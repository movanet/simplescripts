import os
import argparse
import csv
import pdfplumber
import xlrd
from datetime import datetime
from docx import Document
from pptx import Presentation
from typing import List

SUPPORTED_EXTENSIONS = {'.pdf', '.txt', '.docx', '.doc', '.ppt', '.pptx', '.xls', '.xlsx', '.csv'}

def read_txt(file_path: str) -> str:
    with open(file_path, 'r') as file:
        content = file.read()
    return content

def read_docx(file_path: str) -> str:
    doc = Document(file_path)
    content = '\n'.join([para.text for para in doc.paragraphs])
    return content

def read_csv(file_path: str) -> str:
    content = ""
    with open(file_path, 'r') as csvfile:
        reader = csv.reader(csvfile)
        for row in reader:
            content += ' '.join(row) + '\n'
    return content

def read_pdf(file_path: str) -> str:
    content = ""
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            content += page.extract_text() + '\n'
    return content

def read_xls_xlsx(file_path: str) -> str:
    content = ""
    workbook = xlrd.open_workbook(file_path)
    for sheet in workbook.sheets():
        for row in range(sheet.nrows):
            row_values = sheet.row_values(row)
            content += ' '.join(str(value) for value in row_values) + '\n'
    return content

def read_ppt(file_path: str) -> str:
    prs = Presentation(file_path)
    content = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content += shape.text + '\n'
    return content

def process_files(input_folder, output_folder, token_limit):
    all_files = []
    for root, _, files in os.walk(input_folder):
        for file in files:
            all_files.append(os.path.join(root, file))

    now = datetime.now().strftime("%Y%m%d_%H%M%S")
    output_filename = f"{output_folder}/converted_{now}.txt"

    chunks = 1
    with open(output_filename, 'w') as output_file:
        for file_path in all_files:
            file_ext = os.path.splitext(file_path)[1].lower()
            if file_ext not in SUPPORTED_EXTENSIONS or ":Zone.Identifier" in file_path:
                print(f"Error: Unsupported file format: {file_path}")
                continue

            try:
                if file_ext == '.pdf':
                    content = read_pdf(file_path)
                elif file_ext == '.txt':
                    content = read_txt(file_path)
                elif file_ext == '.docx':
                    content = read_docx(file_path)
                elif file_ext == '.doc':
                    content = read_docx(file_path)  # Using docx reader for .doc files as well
                elif file_ext == '.ppt':
                    content = read_ppt(file_path)
                elif file_ext == '.pptx':
                    content = read_ppt(file_path)
                elif file_ext == '.xls':
                    content = read_xls_xlsx(file_path)
                elif file_ext == '.xlsx':
                    content = read_xls_xlsx(file_path)
                elif file_ext == '.csv':
                    content = read_csv(file_path)
                else:
                    print(f"Error: Unsupported file format: {file_path}")
                    continue

                tokens = content.split()
                token_count = len(tokens)

                for i in range(0, token_count, token_limit):
                    chunk_start = i
                    chunk_end = min(i + token_limit, token_count)

                    chunk_content = ' '.join(tokens[chunk_start:chunk_end])

                    if chunk_start == 0:
                        output_file.write(f"### {os.path.basename(file_path)}\n")
                    output_file.write(f"[START CHUNK {chunks}/TOTAL]\n")
                    output_file.write(chunk_content + '\n')
                    output_file.write("##########EndofChunk#########\n")
                    chunks += 1

            except Exception as e:
                print(f"Error processing file {file_path}: {e}")

    print(f"Processed files: {len(all_files)}")
    print(f"Chunks created: {chunks - 1}")
    print(f"Output file: {output_filename}")

if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument('--input_folder', help='Input folder containing the files to process', required=True)
    parser.add_argument('--output_folder', help='Output folder where the text file will be saved', required=True)
    parser.add_argument('--token_limit', type=int, help='Token limit per chunk (positive integer)', required=True)

    args = parser.parse_args()

    if args.token_limit <= 0:
        print("Error: Token limit must be a positive integer.")
        exit(1)

    process_files(args.input_folder, args.output_folder, args.token_limit)

